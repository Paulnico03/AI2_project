from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
UNIGE_BLUE = RGBColor(0, 51, 102)
LIGHT_GRAY = RGBColor(240, 240, 240)
DARK_TEXT = RGBColor(40, 40, 40)

def set_slide_background(slide, color=LIGHT_GRAY):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def format_title(shape, text):
    """Format title"""
    shape.text = text
    shape.text_frame.paragraphs[0].font.size = Pt(44)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = UNIGE_BLUE

def add_bullet_point(text_frame, text, level=0):
    """Add formatted bullet point"""
    if level == 0:
        p = text_frame.paragraphs[0] if text_frame.paragraphs[0].text == "" else text_frame.add_paragraph()
    else:
        p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.space_before = Pt(8)
    p.space_after = Pt(8)

# ============================================================
# SLIDE 1 - TITLE SLIDE
# ============================================================
title_slide_layout = prs.slide_layouts[6]  # Blank layout
slide1 = prs.slides.add_slide(title_slide_layout)
set_slide_background(slide1, UNIGE_BLUE)

# Add title
left = Inches(0.5)
top = Inches(2.5)
width = Inches(9)
height = Inches(1.5)
title_box = slide1.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_frame.text = "Assignment D3-V1"
title_frame.paragraphs[0].font.size = Pt(54)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add subtitle
subtitle_box = slide1.shapes.add_textbox(left, top + Inches(1.2), width, Inches(2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
subtitle_frame.text = "Warehouse Robotics:\nSingle Robot Pick-and-Deliver"
subtitle_frame.paragraphs[0].font.size = Pt(36)
subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add author info
author_box = slide1.shapes.add_textbox(left, top + Inches(3.5), width, Inches(1.5))
author_frame = author_box.text_frame
author_frame.word_wrap = True
author_frame.text = "Paolo Nicolini (s5698969)\nAI for Robotics II — 104731\nUniversità degli Studi di Genova"
author_frame.paragraphs[0].font.size = Pt(16)
author_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 2 - PROBLEM & DOMAIN OVERVIEW (~2 min)
# ============================================================
blank_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)

# Title
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame2 = title_box2.text_frame
title_frame2.text = "Problem & Domain Overview"
title_frame2.paragraphs[0].font.size = Pt(40)
title_frame2.paragraphs[0].font.bold = True
title_frame2.paragraphs[0].font.color.rgb = UNIGE_BLUE

# Content
content_box2 = slide2.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.5))
tf2 = content_box2.text_frame
tf2.word_wrap = True

add_bullet_point(tf2, "Scenario: A mobile robot navigates a warehouse graph of locations (dock, storage aisles, shipping bay)")
add_bullet_point(tf2, "Critical Constraint: The robot can carry only ONE package at a time")
add_bullet_point(tf2, "Goal: Deliver all packages from storage to shipping dock")
add_bullet_point(tf2, "Challenge: This capacity constraint forces sequential, back-and-forth planning patterns")

# Speaker note
note_box2 = slide2.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.8))
note_frame2 = note_box2.text_frame
note_frame2.text = "🗣️ Keep it intuitive: Draw the warehouse graph mentally. Emphasize the bottleneck of single-capacity."
note_frame2.paragraphs[0].font.size = Pt(12)
note_frame2.paragraphs[0].font.italic = True
note_frame2.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ============================================================
# SLIDE 3 - CLASSICAL PDDL MODEL DESIGN CHOICES (~1.5 min)
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)

title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame3 = title_box3.text_frame
title_frame3.text = "Classical PDDL: Discrete Domain Modelling"
title_frame3.paragraphs[0].font.size = Pt(40)
title_frame3.paragraphs[0].font.bold = True
title_frame3.paragraphs[0].font.color.rgb = UNIGE_BLUE

content_box3 = slide3.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.3))
tf3 = content_box3.text_frame
tf3.word_wrap = True

add_bullet_point(tf3, "Key Predicates: (hand-empty), (holding), (at-robot), (at-package)")
add_bullet_point(tf3, "Actions: move, pick, drop — each is explicit and separate")
add_bullet_point(tf3, "Capacity Enforcement: (pick) requires (hand-empty). This forces sequential execution.")
add_bullet_point(tf3, "Result: A 3-package problem produces a 12-step back-and-forth plan")

note_box3 = slide3.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.8))
note_frame3 = note_box3.text_frame
note_frame3.text = "🗣️ Walk through the pick action: show how (hand-empty) blocks multiple pickups. This is the design choice."
note_frame3.paragraphs[0].font.size = Pt(12)
note_frame3.paragraphs[0].font.italic = True
note_frame3.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ============================================================
# SLIDE 4 - PDDL+ CONTINUOUS EXTENSION (~1.5 min)
# ============================================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)

title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame4 = title_box4.text_frame
title_frame4.text = "PDDL+ Extension: Hybrid Discrete-Continuous"
title_frame4.paragraphs[0].font.size = Pt(40)
title_frame4.paragraphs[0].font.bold = True
title_frame4.paragraphs[0].font.color.rgb = UNIGE_BLUE

content_box4 = slide4.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.3))
tf4 = content_box4.text_frame
tf4.word_wrap = True

add_bullet_point(tf4, "Continuous Processes: robot-transit (distance counts down) and package-aging (deadline counts down)")
add_bullet_point(tf4, "Numeric Fluents: (distance-left ?r), (time-remaining ?p) — these evolve continuously over time")
add_bullet_point(tf4, "Automatic Events: deadline-breach fires when time hits zero, marking the package as 'missed-deadline'")
add_bullet_point(tf4, "Critical Fix: Event precondition (not (missed-deadline ?p)) prevents infinite firing — ensures it fires ONCE")

note_box4 = slide4.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.8))
note_frame4 = note_box4.text_frame
note_frame4.text = "🗣️ Emphasize the event design. The third precondition stops infinite loops — this is the key insight."
note_frame4.paragraphs[0].font.size = Pt(12)
note_frame4.paragraphs[0].font.italic = True
note_frame4.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ============================================================
# SLIDE 5 - LIVE DEMO / PLAN OUTPUT (~3 min)
# ============================================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)

title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame5 = title_box5.text_frame
title_frame5.text = "Live Demo: ENHSP Planner Output"
title_frame5.paragraphs[0].font.size = Pt(40)
title_frame5.paragraphs[0].font.bold = True
title_frame5.paragraphs[0].font.color.rgb = UNIGE_BLUE

# Add code snippet (simplified plan)
code_box5 = slide5.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(4.5))
code_frame5 = code_box5.text_frame
code_frame5.word_wrap = True

code_text = """Classical Problem 2 (3-package delivery):
0.0:  (move robby dock aisle-A)
1.0:  (pick robby pkg1 aisle-A)
2.0:  (move robby aisle-A shipping)
3.0:  (drop robby pkg1 shipping)
4.0:  (move robby shipping aisle-B)
...
11.0: (drop robby pkg2 shipping)

Result: 12 steps total. Back-and-forth pattern visible.

PDDL+ Problem 2 (same scenario, with time):
0:    (start-move robby dock aisle-A)
0:    -----waiting---- [2.0]
2.0:  (end-move robby dock aisle-A)
2.0:  (pick robby pkg2 aisle-A)
...
14.0: (deliver-package robby pkg1 shipping)

Result: All packages delivered. Makespan = 14 time units."""

code_frame5.text = code_text
code_frame5.paragraphs[0].font.size = Pt(11)
code_frame5.paragraphs[0].font.name = 'Courier New'
code_frame5.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

note_box5 = slide5.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(8.6), Inches(1.2))
note_frame5 = note_box5.text_frame
note_frame5.word_wrap = True
note_frame5.text = "🗣️ Open VS Code. Show the terminal output. Walk through the plan step-by-step. Highlight the waiting blocks in PDDL+ (showing continuous time)."
note_frame5.paragraphs[0].font.size = Pt(11)
note_frame5.paragraphs[0].font.italic = True
note_frame5.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ============================================================
# SLIDE 6 - REFLECTION & CONCLUSIONS (~2 min)
# ============================================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)

title_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame6 = title_box6.text_frame
title_frame6.text = "Reflection & Key Takeaways"
title_frame6.paragraphs[0].font.size = Pt(40)
title_frame6.paragraphs[0].font.bold = True
title_frame6.paragraphs[0].font.color.rgb = UNIGE_BLUE

content_box6 = slide6.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.3))
tf6 = content_box6.text_frame
tf6.word_wrap = True

add_bullet_point(tf6, "Main Challenge: Exponential state-space growth when adding packages in continuous time (1700× nodes for 3 packages)")
add_bullet_point(tf6, "Heuristic Pitfall: Delete-relaxation heuristics miss deadline constraints. Solution: explicitly embed failure in goal.")
add_bullet_point(tf6, "Model Limitation: Single-capacity severely limits throughput. Future: multi-gripper models needed for realism.")
add_bullet_point(tf6, "Key Insight: PDDL+ cleanly models time-dependent constraints; careful numerical tuning is essential for scalability.")

note_box6 = slide6.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.8))
note_frame6 = note_box6.text_frame
note_frame6.text = "🗣️ Summarize what you learned. Mention what you'd do differently. Be honest about limitations. Thank you & Q&A."
note_frame6.paragraphs[0].font.size = Pt(11)
note_frame6.paragraphs[0].font.italic = True
note_frame6.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

# ============================================================
# SAVE PRESENTATION
# ============================================================
prs.save('presentation.pptx')
print("✅ Presentation generated successfully: presentation.pptx")
print("📊 Total slides: 6 (1 title + 5 content)")
print("⏱️  Timing: ~2 + ~3 + ~3 + ~2 = ~10 minutes")